#!/usr/bin/env python3
"""
ChordPro → PPTX lyric slide generator

Generates a 16:9 PPTX with a black background and a single centered text box
in the top-third of the slide. Each slide contains up to two lyric lines.

Usage:
  - Single file → PPTX:
      python3 scripts/chordpro_to_pptx.py path/to/song.chordpro -o out.pptx

  - Batch (directory → PPTX files into output folder):
      python3 scripts/chordpro_to_pptx.py path/to/dir -O out_pptx/

Requirements:
  pip install python-pptx

Notes:
  - Ignores metadata/directives (e.g., {title: ...}), section headers ([Verse], [Chorus]),
    and chord symbols ([C], [G/D]) embedded in lyric lines.
  - Uses bold Calibri, white, center-aligned, 58pt (configurable).
  - Text box is horizontally centered; vertically placed in the top third.
"""

from __future__ import annotations

import argparse
import os
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.dml import MSO_THEME_COLOR
except ImportError as e:
    print("This script requires python-pptx. Install with: pip install python-pptx", file=sys.stderr)
    raise


CHORD_INLINE_PATTERN = re.compile(r"\[[^\]]+\]")
DIRECTIVE_LINE = re.compile(r"^\s*\{[^}]*\}\s*$")
SECTION_HEADER = re.compile(r"^\s*\[[^\]]+\]\s*$")
COMMENT_LINE = re.compile(r"^\s*#")


def parse_chordpro_lyrics(path: Path) -> list[str]:
    """Parse a .chordpro file and return plain lyric lines.

    - Drops directive lines like {title: ...}, {comment: ...}
    - Drops section headers like [Verse], [Chorus]
    - Strips inline chord symbols like [C], [G/D]
    - Returns non-empty lyric lines, trimmed of leading/trailing whitespace.
    """
    lines: list[str] = []
    text = path.read_text(encoding="utf-8", errors="replace")
    for raw in text.splitlines():
        if not raw.strip():
            # Keep structure simple: treat blank as separator (ignored for slides)
            continue
        if DIRECTIVE_LINE.match(raw):
            continue
        if SECTION_HEADER.match(raw):
            continue
        if COMMENT_LINE.match(raw):
            continue

        # Remove inline chord tags like [C], [G/D]
        lyric = CHORD_INLINE_PATTERN.sub("", raw)
        lyric = lyric.strip()
        if lyric:
            lines.append(lyric)
    return lines


def chunk_lines(lines: list[str], n: int = 2, pad_last: bool = False) -> list[list[str]]:
    """Group lines into chunks of n (default 2). Optionally pad the last chunk."""
    chunks: list[list[str]] = []
    buf: list[str] = []
    for ln in lines:
        buf.append(ln)
        if len(buf) == n:
            chunks.append(buf)
            buf = []
    if buf:
        if pad_last and len(buf) < n:
            buf += [""] * (n - len(buf))
        chunks.append(buf)
    return chunks


def make_presentation(
    lyric_pairs: list[list[str]],
    font_name: str = "Calibri",
    font_size_pt: float = 58.0,
    bold: bool = True,
    slide_width_in: float = 13.333,  # 16:9
    slide_height_in: float = 7.5,
    box_width_ratio: float = 0.9,
    box_height_ratio: float = 0.28,
    box_center_y_ratio: float = 1 / 6,  # center of box at top third center
) -> Presentation:
    """Create a Presentation object with the given lyric pairs (two lines per slide)."""
    prs = Presentation()
    prs.slide_width = Inches(slide_width_in)
    prs.slide_height = Inches(slide_height_in)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    box_w = int(slide_w * box_width_ratio)
    box_h = int(slide_h * box_height_ratio)
    left = int((slide_w - box_w) / 2)
    top = int(slide_h * box_center_y_ratio - box_h / 2)

    for pair in lyric_pairs:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        # Black background
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)

        # Text box
        shape = slide.shapes.add_textbox(left, top, box_w, box_h)
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True

        # Build paragraphs
        for i, line in enumerate(pair):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.alignment = PP_ALIGN.CENTER
            # Style the run(s) for this paragraph
            if p.runs:
                # python-pptx creates one run by default for p.text
                r = p.runs[0]
            else:
                r = p.add_run()
            r.font.name = font_name
            r.font.size = Pt(font_size_pt)
            r.font.bold = bold
            r.font.color.rgb = RGBColor(255, 255, 255)
            # Ensure theme does not override to auto
            r.font.color.theme_color = None

    return prs


def ensure_outdir(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True)


def main(argv: list[str] | None = None) -> int:
    ap = argparse.ArgumentParser(description="Generate a PPTX from ChordPro lyrics (two lines per slide).")
    ap.add_argument("input", help="Path to a .chordpro file or a directory containing .chordpro files.")
    ap.add_argument("-o", "--output", help="Output PPTX path (for single file input). Defaults to <input_basename>.pptx.")
    ap.add_argument("-O", "--output-dir", help="Output directory for batch mode (when input is a directory). Defaults to ./pptx_out.")
    ap.add_argument("--font-name", default="Calibri", help="Font name (default: Calibri)")
    ap.add_argument("--font-size", type=float, default=58.0, help="Font size in points (default: 58)")
    ap.add_argument("--pad-last", action="store_true", help="Pad last slide to exactly two lines (adds a blank line if needed).")

    args = ap.parse_args(argv)
    src = Path(args.input)

    if not src.exists():
        print(f"Input not found: {src}", file=sys.stderr)
        return 1

    if src.is_file():
        if src.suffix.lower() not in {".chordpro", ".cho", ".crd"}:
            print("Warning: input does not look like a ChordPro file (extensions: .chordpro, .cho, .crd)", file=sys.stderr)
        lines = parse_chordpro_lyrics(src)
        pairs = chunk_lines(lines, n=2, pad_last=args.pad_last)
        prs = make_presentation(
            pairs,
            font_name=args.font_name,
            font_size_pt=args.font_size,
        )
        out = Path(args.output) if args.output else src.with_suffix(".pptx")
        prs.save(out)
        print(f"Wrote {out}")
        return 0

    # Directory mode: generate per file into output-dir
    outdir = Path(args.output_dir) if args.output_dir else Path("pptx_out")
    ensure_outdir(outdir)
    count = 0
    for path in sorted(src.rglob("*")):
        if not path.is_file():
            continue
        if path.suffix.lower() not in {".chordpro", ".cho", ".crd"}:
            continue
        try:
            lines = parse_chordpro_lyrics(path)
            pairs = chunk_lines(lines, n=2, pad_last=args.pad_last)
            prs = make_presentation(
                pairs,
                font_name=args.font_name,
                font_size_pt=args.font_size,
            )
            out = outdir / (path.stem + ".pptx")
            prs.save(out)
            print(f"Wrote {out}")
            count += 1
        except Exception as e:
            print(f"Failed for {path}: {e}", file=sys.stderr)
    if count == 0:
        print("No .chordpro files found.")
    else:
        print(f"Done. Generated {count} PPTX files in {outdir}.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

